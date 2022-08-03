from gtts import gTTS
import streamlit as st
import pdfplumber
import docx2txt
import pandas as pd
import time
import pip
from pptx import Presentation
import easyocr
import aspose.words as aw
from englisttohindi.englisttohindi import EngtoHindi




# pip.main(["install", "openpyxl"])
text = []
hind = []
a = ""
counter=0
count = 0
n=0
total = """"""
a = []
def convert_to_audio(text):
    global counter
    print(text)
    counter +=1
    audio = gTTS(text=text, lang='en', tld='com')
    audio.save(f"textaud{counter}.mp3")
    audio_file = open(f"textaud{counter}.mp3", "rb")
    st.audio(audio_file.read())
    with open(f"textaud{counter}.mp3", 'rb') as f:
        st.download_button('Download Mp3', f, file_name=f"textaud{counter}.mp3")


def convert_to_hindiaudio(text):
    global counter
    counter +=1
    audio = gTTS(text=text, lang='hi', tld='com')
    audio.save(f"textaud{counter}.mp3")
    audio_file = open(f"textaud{counter}.mp3", "rb")
    st.audio(audio_file.read())
    with open(f"textaud{counter}.mp3", 'rb') as f:
        st.download_button('Download Mp3', f, file_name=f"textaud{counter}.mp3")


count = 0
interface = ""
total_text = """"""
raw_text = ""
menu = ["Scanned PDF","PDF File","DocumentFiles","ExcelFile","PPT"]

choice = st.sidebar.selectbox("Menu",menu)

if choice == "PDF File":
    st.title("PDF File")
    docx_file = st.file_uploader("Upload Document", type=["pdf"])
    if docx_file is not None:
        interface = "File has been uploaded"
        st.subheader(interface)
        if docx_file.type == "application/pdf":
            starting_page_no = st.number_input('enter the starting page number of document', value=1)
            ending_page_no = st.number_input('enter the ending page number of document', value=1)
            total_pages = ending_page_no - starting_page_no
            value = ['english','hindi']
            stm = st.selectbox('Select Language',value)
            if st.button("Process"):
                if stm == 'english':
                    # file_details = {"filename": docx_file.name, "filetype": docx_file.type,
                    #                 "filesize": docx_file.size}
                    # st.write(file_details)
                    for i in range(total_pages + 1):
                        loop = 1
                        try:
                            with pdfplumber.open(docx_file) as pdf:
                                pages = pdf.pages[int(starting_page_no - 1 + i)]
                                line = pages.extract_text()
                                total_text = total_text + line
                                if len(total_text) < 40000:
                                    pass
                                else:
                                    convert_to_audio(total_text)
                                    time.sleep(999)
                                    total_text = ""
                                    loop = 0

                        except:
                            st.write("None")
                        # time.sleep(2)
                    if loop == 1 and len(total_text)>0:
                        print(len(total_text))
                        convert_to_audio(total_text)
                    else:
                        print(docx_file.name)
                        doc = aw.Document(docx_file.name)
                        for page in range(0, doc.page_count):
                            n += 1
                            extractedPage = doc.extract_pages(page, 1)
                            #         save the image file in output folder
                            extractedPage.save(f"Output_{page + n}.jpg")
                            break
                        reader = easyocr.Reader(['en'])
                        text = reader.readtext('Output_1.jpg')
                        text = text[3:]
                        for i in text:
                            a.append(i[1])
                        b = " ".join(a)
                        print(b)
                        convert_to_audio(b)

                else:
                    # print('hindi')
                    # file_details = {"filename": docx_file.name, "filetype": docx_file.type,
                    #                 "filesize": docx_file.size}
                    # st.write(file_details)
                    for i in range(total_pages + 1):

                        loop = 1
                        try:
                            with pdfplumber.open(docx_file) as pdf:
                                pages = pdf.pages[int(starting_page_no - 1 + i)]
                                line = pages.extract_text()
                                total_text = total_text + line
                                # if total_text is None:
                                #     total_text = f'no text in {count} page'
                                if len(total_text) < 40000:
                                    pass
                                else:
                                    convert_to_audio(total_text)
                                    print("very much texts")
                                    time.sleep(999)
                                    total_text = ""
                                    loop = 0

                        except:
                            st.write("None")
                        time.sleep(2)
                    if loop == 1:
                        # print(total_text)
                        if total_text is None or len(total_text)==0:
                            total_text = f'no word in this page'
                        s = total_text.split(' ')
                        for i in s:
                            if len(i)>2:
                                a.append(i)
                        s = " ".join(a)
                        x = s.replace('\n',' ')
                        c= x.lower()
                        for i in c:
                            text = c.split(',')
                        # print(text)
                        for i in text:
                            trans = EngtoHindi(message=i)
                            j = trans.convert
                            hind.append(j)
                            # print(hind)
                        a = "".join(hind)
                        # print(a)

                        convert_to_hindiaudio(a)

                    # st.write("If you last converted audio file were more than 1 hour than please wait for 10 minute")
                    # st.write(len(total_text))
if choice == "Scanned PDF":
    st.title('Scanned PDF')
    docx_file = st.file_uploader("Upload Document", type=["pdf"])
    if docx_file is not None:
        doc = aw.Document(docx_file)
        starting_page_no = st.number_input('enter the starting page number of document', value=1)
        ending_page_no = st.number_input('enter the ending page number of document', value=1)
        total_pages = ending_page_no - starting_page_no
        value = ['english','hindi']
        stm = st.selectbox('Select Language', value)
        if st.button("Process"):
            if stm == 'english':
                for i in range(total_pages+1):
                    extractedPage = doc.extract_pages(int(starting_page_no - 1 + i),1)
                    a= extractedPage.save(f"Output_{i}.jpg")
                    reader = easyocr.Reader(['en'])
                    text = reader.readtext(f"Output_{i}.jpg")
                for i in text:
                    count+=1
                    if count>1 and count< (len(text) -2):
                        total_text+=i[1]
                # print(total_text)
                convert_to_audio(total_text)
            else:
                for i in range(total_pages+1):
                    extractedPage = doc.extract_pages(int(starting_page_no - 1 + i),1)
                    a= extractedPage.save(f"Output_{i}.jpg")
                    reader = easyocr.Reader(['en'])
                    text = reader.readtext(a)
                for i in text:
                    count+=1
                    if count>1 and count< (len(text) -2):
                        total_text+=i[1]

                s = total_text.split(' ')
                for i in s:
                    if len(i) > 2:
                        a.append(i)
                s = " ".join(a)
                x = s.replace('\n', '')
                c = x.lower()
                # print(c)
                trans = EngtoHindi(message=c)
                j = trans.convert
                # print(j)
                convert_to_hindiaudio(j)
if choice == "DocumentFiles":
    st.title("DocumentFiles")
    docx_file = st.file_uploader("Upload Document", type=["docx", "txt"])
    if docx_file is not None:
        interface = "File has been uploaded"
        st.subheader(interface)
        if docx_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            value = ['english','hindi']
            stm = st.selectbox('Select Language', value)
            if st.button("Process"):
                if stm == 'english':
                    # file_details = {"filename": docx_file.name, "filetype": docx_file.type,
                    #                 "filesize": docx_file.size}
                    # st.write(file_details)
                    total_text = docx2txt.process(docx_file)
                    convert_to_audio(total_text)
                    # audio_file = open("textaud.mp3", "rb")
                    # st.audio(audio_file.read())
                else:
                    total_text = docx2txt.process(docx_file)
                    s = total_text.split(' ')
                    for i in s:
                        if len(i) > 2:
                            a.append(i)
                    s = " ".join(a)
                    x = s.replace('\n', '')
                    c = x.lower()
                    # print(c)
                    trans = EngtoHindi(message=c)
                    j = trans.convert
                    # print(j)
                    convert_to_hindiaudio(j)
        elif docx_file.type == "text/plain":
            value = ['english','hindi']
            stm = st.selectbox('Select Language', value)
            if st.button("Process"):
                if stm == 'english':
                    # file_details = {"filename": docx_file.name, "filetype": docx_file.type,
                    #                 "filesize": docx_file.size}
                    # st.write(file_details)
                    total_text = str(docx_file.read(), "utf-8")
                    convert_to_audio(total_text)
                else:
                    total_text = str(docx_file.read(), "utf-8")
                    s = total_text.split(' ')
                    for i in s:
                        if len(i) > 2:
                            a.append(i)
                    s = " ".join(a)
                    x = s.replace('\n', '')
                    c = x.lower()
                    # print(c)
                    trans = EngtoHindi(message=c)
                    j = trans.convert
                    # print(j)
                    convert_to_hindiaudio(j)
        elif docx_file is None:
            st.subheader(interface)
    else:
        st.subheader(interface)
if choice == "ExcelFile":
    st.subheader("Dataset")
    data_file = st.file_uploader("Upload CSV", type=["xlsx"])
    if data_file is not None:
        interface = "file has been uploaded"
        st.subheader(interface)
        value = ['english','hindi']
        stm = st.selectbox('Select Language', value)
        if st.button("Process"):
            if stm =='english':
                # file_details = {"filename": data_file.name, "filetype": data_file.type,
                #                 "filesize": data_file.size}
                xl = pd.ExcelFile(data_file)
                for sheet in xl.sheet_names:
                    file = pd.read_excel(xl, sheet_name=sheet)
                    docx_file = file.to_csv(sheet + '.txt', header=False, index=False)
                    docx_file = open('Sheet1.txt','rb')
                    total_text = str(docx_file.read(), "utf-8")
                    convert_to_audio(total_text)
            else:
                xl = pd.ExcelFile(data_file)
                for sheet in xl.sheet_names:
                    file = pd.read_excel(xl, sheet_name=sheet)
                    docx_file = file.to_csv(sheet + '.txt', header=False, index=False)
                    docx_file = open('Sheet1.txt', 'rb')
                    total_text = str(docx_file.read(), "utf-8")
                    # convert_to_audio(total_text)
                s = total_text.split(' ')
                for i in s:
                    if len(i) > 2:
                        a.append(i)
                s = " ".join(a)
                x = s.replace('\n', '')
                c = x.lower()
                # print(c)
                trans = EngtoHindi(message=c)
                j = trans.convert
                # print(j)
                convert_to_hindiaudio(j)
    else:
        st.subheader(interface)

if choice == "PPT":
    st.title("PPT")
    data_file = st.file_uploader("Upload CSV", type=["pptx"])
    if data_file is not None:
        interface = "file has been uploaded"
        st.subheader(interface)
        # file_details = {"filename": data_file.name, "filetype": data_file.type,
        #                 "filesize": data_file.size}
        final_text = ""
        # for eachfile in glob.glob("*pptx"):
        prs = Presentation(data_file)
        print(data_file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    total_text += shape.text
        print(total_text)
        value = ['english','hindi']
        stm = st.selectbox('Select Language', value)
        if st.button("Process"):
            if stm == 'english':
                convert_to_audio(total_text)
            else:
                s = total_text.split(' ')
                for i in s:
                    if len(i) > 2:
                        a.append(i)
                s = " ".join(a)
                x = s.replace('\n', '')
                c = x.lower()
                # print(c)
                trans = EngtoHindi(message=c)
                j = trans.convert
                # print(j)
                convert_to_hindiaudio(j)
    else:
        st.subheader(interface)
